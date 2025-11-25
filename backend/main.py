from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, JSONResponse
import shutil, subprocess, uuid, os, pathlib, tempfile, zipfile

app = FastAPI(title="PDF Tools API", version="1.0.0")

TMP = '/tmp/pdf_tools'
os.makedirs(TMP, exist_ok=True)

def save_upload(upload: UploadFile):
    fn = os.path.join(TMP, f"{uuid.uuid4().hex}_{upload.filename}")
    with open(fn, 'wb') as f:
        shutil.copyfileobj(upload.file, f)
    return fn

def run_cmd(cmd):
    proc = subprocess.run(cmd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    return proc.returncode, proc.stdout.decode('utf-8', errors='ignore'), proc.stderr.decode('utf-8', errors='ignore')

@app.post('/convert/pdf-to-word')
async def pdf_to_word(file: UploadFile = File(...)):
    input_path = save_upload(file)
    outdir = os.path.dirname(input_path)
    cmd = f"soffice --headless --convert-to docx --outdir {outdir} {input_path}"
    rc, outp, err = run_cmd(cmd)
    if rc != 0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':outp,'err':err})
    # find produced docx
    docs = list(pathlib.Path(outdir).glob('*.docx'))
    if docs:
        docs_sorted = sorted(docs, key=lambda p: p.stat().st_mtime)
        return FileResponse(str(docs_sorted[-1]), filename=docs_sorted[-1].name)
    raise HTTPException(status_code=500, detail='docx not produced')

@app.post('/convert/word-to-pdf')
async def word_to_pdf(file: UploadFile = File(...)):
    input_path = save_upload(file)
    outdir = os.path.dirname(input_path)
    cmd = f"soffice --headless --convert-to pdf --outdir {outdir} {input_path}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    pdfs = list(pathlib.Path(outdir).glob('*.pdf'))
    if pdfs:
        return FileResponse(str(sorted(pdfs, key=lambda p: p.stat().st_mtime)[-1]), filename=sorted(pdfs)[-1].name)
    raise HTTPException(status_code=500, detail='pdf not produced')

@app.post('/convert/pdf-to-jpg')
async def pdf_to_jpg(file: UploadFile = File(...), dpi: int = Form(150)):
    input_path = save_upload(file)
    out_prefix = os.path.join(TMP, uuid.uuid4().hex + '_page')
    cmd = f"pdftoppm -jpeg -r {dpi} {input_path} {out_prefix}"

    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    imgs = sorted([str(p) for p in pathlib.Path(TMP).glob(out_prefix.split('/')[-1]+'*')])
    # fallback gather
    if not imgs:
        imgs = sorted([str(p) for p in pathlib.Path(TMP).glob('*_page*')])
    return JSONResponse({'images': imgs})

@app.post('/convert/jpg-to-pdf')
async def jpg_to_pdf(file: UploadFile = File(...)):
    input_path = save_upload(file)
    out_path = os.path.join(TMP, uuid.uuid4().hex + '.pdf')
    cmd = f"img2pdf {input_path} -o {out_path}"

    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out_path, filename=os.path.basename(out_path))

@app.post('/tools/merge')
async def merge_pdfs(files: list[UploadFile] = File(...)):
    saved = []
    for f in files:
        saved.append(save_upload(f))
    out = os.path.join(TMP, uuid.uuid4().hex + '_merged.pdf')
    # try pdfunite then ghostscript fallback
    cmd = f"pdfunite {' '.join(saved)} {out} || gs -q -dNOPAUSE -sDEVICE=pdfwrite -sOUTPUTFILE={out} -dBATCH {' '.join(saved)}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out, filename='merged.pdf')

@app.post('/tools/split')
async def split_pdf(file: UploadFile = File(...), ranges: str = Form(...)):
    input_path = save_upload(file)
    outdir = os.path.join(TMP, uuid.uuid4().hex + '_split')
    os.makedirs(outdir, exist_ok=True)
    try:
        from PyPDF2 import PdfReader, PdfWriter
        reader = PdfReader(input_path)
        parts = []
        def pages_from_range(r):
            if '-' in r:
                a,b = r.split('-',1); return range(int(a)-1,int(b))
            return [int(r)-1]
        i=0
        for seg in ranges.split(','):
            pages = pages_from_range(seg)
            writer = PdfWriter()
            for p in pages:
                writer.add_page(reader.pages[p])
            outp = os.path.join(outdir, f'part_{i+1}.pdf')
            with open(outp,'wb') as f: writer.write(f)
            parts.append(outp); i+=1
        zip_path = outdir + '.zip'
        with zipfile.ZipFile(zip_path,'w') as zf:
            for p in parts: zf.write(p, os.path.basename(p))
        return FileResponse(zip_path, filename='split_parts.zip')
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post('/tools/ocr')
async def ocr_pdf(file: UploadFile = File(...), lang: str = Form('eng')):
    input_path = save_upload(file)
    try:
        from pdf2image import convert_from_path
        import pytesseract
        pages = convert_from_path(input_path, dpi=300)
        text_all = []
        for p in pages:
            text = pytesseract.image_to_string(p, lang=lang)
            text_all.append(text)
        text_out = '\n\n'.join(text_all)
        outfile = input_path + '.txt'
        with open(outfile,'w', encoding='utf-8') as f: f.write(text_out)
        return FileResponse(outfile, filename=os.path.basename(outfile))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get('/health')
async def health():
    return {'status':'ok'}



# === Additional endpoints implemented (watermark, page numbers, compress, protect/unlock, rotate, reorder, extract, delete, pdf->excel, excel->pdf, ppt->pdf, pdf->ppt, repair, pdfa, html->pdf, edit text) ===
from fastapi import UploadFile, File, Form
from fastapi.responses import FileResponse, JSONResponse
import pathlib, subprocess, uuid, os, zipfile
from PyPDF2 import PdfReader, PdfWriter
import io

# Helper: create temp output path
def make_out_path(suffix):
    return os.path.join(TMP, f"{uuid.uuid4().hex}{suffix}")

# Watermark (text) using reportlab to produce overlay PDF then merge
@app.post('/tools/watermark-text')
async def watermark_text(file: UploadFile = File(...), text: str = Form(...), opacity: float = Form(0.15), fontsize: int = Form(40)):
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from PyPDF2 import PdfReader, PdfWriter
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing libraries: {e}")

    input_path = save_upload(file)
    reader = PdfReader(input_path)
    out_pdf = make_out_path('_watermarked.pdf')

    # Create watermark PDF with reportlab for each page size: we will create one large page and scale later
    wm_path = make_out_path('_wm.pdf')
    c = canvas.Canvas(wm_path, pagesize=letter)
    c.setFillAlpha(opacity)
    c.setFont("Helvetica", fontsize)
    c.translate(300, 400)
    c.rotate(45)
    c.drawString(0, 0, text)
    c.save()

    # Merge watermark onto each page
    writer = PdfWriter()
    wm_reader = PdfReader(wm_path)
    for p in reader.pages:
        p.merge_page(wm_reader.pages[0])
        writer.add_page(p)
    with open(out_pdf, 'wb') as f:
        writer.write(f)
    # cleanup watermark tmp
    try:
        os.remove(wm_path)
    except: pass
    return FileResponse(out_pdf, filename='watermarked.pdf')

# Page numbers
@app.post('/tools/add-page-numbers')
async def add_page_numbers(file: UploadFile = File(...), fmt: str = Form("Page {n} of {N}"), position: str = Form("bottom")):
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from PyPDF2 import PdfReader, PdfWriter
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing libraries: {e}")

    input_path = save_upload(file)
    reader = PdfReader(input_path)
    total = len(reader.pages)
    out_pdf = make_out_path('_pagenums.pdf')

    # create watermark per page size individually and merge
    writer = PdfWriter()
    for idx, page in enumerate(reader.pages):
        # create single-page watermark
        wm_path = make_out_path(f'_pnum_{idx}.pdf')
        from reportlab.lib.units import mm
        from reportlab.pdfgen import canvas
        w = page.mediabox.width
        h = page.mediabox.height
        c = canvas.Canvas(wm_path, pagesize=(w, h))
        text = fmt.replace('{n}', str(idx+1)).replace('{N}', str(total))
        if position == 'bottom':
            x = w/2 - 50; y = 20
        else:
            x = w/2 - 50; y = h - 30
        c.setFont("Helvetica", 10)
        c.drawString(x, y, text)
        c.save()
        wm_reader = PdfReader(wm_path)
        page.merge_page(wm_reader.pages[0])
        writer.add_page(page)
        try:
            os.remove(wm_path)
        except: pass
    with open(out_pdf, 'wb') as f:
        writer.write(f)
    return FileResponse(out_pdf, filename='pagenums.pdf')

# Compress PDF using Ghostscript presets
@app.post('/tools/compress')
async def compress_pdf(file: UploadFile = File(...), preset: str = Form("screen")):
    # presets: screen, ebook, printer, prepress
    preset_map = {"high":"/printer", "regular":"/ebook", "extreme":"/screen", "screen":"/screen", "ebook":"/ebook", "printer":"/printer"}
    input_path = save_upload(file)
    out = make_out_path('_compressed.pdf')
    gs_preset = preset_map.get(preset, "/screen")
    cmd = f"gs -sDEVICE=pdfwrite -dCompatibilityLevel=1.4 -dPDFSETTINGS={gs_preset} -dNOPAUSE -dQUIET -dBATCH -sOutputFile={out} {input_path}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out, filename='compressed.pdf')

# Protect PDF (encrypt)
@app.post('/tools/protect')
async def protect_pdf(file: UploadFile = File(...), user_pass: str = Form(...), owner_pass: str = Form(...)):
    input_path = save_upload(file)
    out = make_out_path('_protected.pdf')
    cmd = f"qpdf --encrypt {user_pass} {owner_pass} 256 -- {input_path} {out}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out, filename='protected.pdf')

# Unlock PDF (decrypt) - requires password
@app.post('/tools/unlock')
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    input_path = save_upload(file)
    out = make_out_path('_unlocked.pdf')
    cmd = f"qpdf --password={password} --decrypt {input_path} {out}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out, filename='unlocked.pdf')

# Rotate pages
@app.post('/tools/rotate')
async def rotate_pages(file: UploadFile = File(...), pages: str = Form(...), angle: int = Form(90)):
    # pages: "1,3,5-7" rotate pages by angle
    input_path = save_upload(file)
    reader = PdfReader(input_path)
    writer = PdfWriter()
    def resolve_pages(spec):
        out = set()
        for part in spec.split(','):
            if '-' in part:
                a,b = part.split('-',1); out.update(range(int(a)-1,int(b)))
            else:
                out.add(int(part)-1)
        return out
    to_rotate = resolve_pages(pages)
    for i,p in enumerate(reader.pages):
        if i in to_rotate:
            p.rotate(angle)
        writer.add_page(p)
    outp = make_out_path('_rotated.pdf')
    with open(outp,'wb') as f: writer.write(f)
    return FileResponse(outp, filename='rotated.pdf')

# Reorder pages - provide order like "3,1,2,4"
@app.post('/tools/reorder')
async def reorder_pages(file: UploadFile = File(...), order: str = Form(...)):
    input_path = save_upload(file)
    reader = PdfReader(input_path)
    writer = PdfWriter()
    indices = [int(x)-1 for x in order.split(',')]
    for idx in indices:
        writer.add_page(reader.pages[idx])
    outp = make_out_path('_reordered.pdf')
    with open(outp,'wb') as f: writer.write(f)
    return FileResponse(outp, filename='reordered.pdf')

# Extract pages (save specified pages as new PDF)
@app.post('/tools/extract')
async def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):
    input_path = save_upload(file)
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for pnum in pages.split(','):
        if '-' in pnum:
            a,b = pnum.split('-',1)
            for i in range(int(a)-1,int(b)):
                writer.add_page(reader.pages[i])
        else:
            writer.add_page(reader.pages[int(pnum)-1])
    outp = make_out_path('_extracted.pdf')
    with open(outp,'wb') as f: writer.write(f)
    return FileResponse(outp, filename='extracted.pdf')

# Delete pages - create new PDF without specified pages
@app.post('/tools/delete-pages')
async def delete_pages(file: UploadFile = File(...), pages: str = Form(...)):
    input_path = save_upload(file)
    reader = PdfReader(input_path)
    writer = PdfWriter()
    del_set = set()
    for pnum in pages.split(','):
        if '-' in pnum:
            a,b = pnum.split('-',1)
            del_set.update(range(int(a)-1,int(b)))
        else:
            del_set.add(int(pnum)-1)
    for i,p in enumerate(reader.pages):
        if i not in del_set:
            writer.add_page(p)
    outp = make_out_path('_deleted.pdf')
    with open(outp,'wb') as f: writer.write(f)
    return FileResponse(outp, filename='cleaned.pdf')

# PDF -> Excel (using camelot)
@app.post('/convert/pdf-to-excel')
async def pdf_to_excel(file: UploadFile = File(...)):
    input_path = save_upload(file)
    try:
        import camelot
        import pandas as pd
        tables = camelot.read_pdf(input_path, pages='all')
        if tables.n == 0:
            raise Exception('No tables found')
        out_xlsx = make_out_path('.xlsx')
        # concatenate tables into sheets
        import openpyxl
        from openpyxl import Workbook
        wb = Workbook()
        wb.remove(wb.active)
        for i, t in enumerate(tables):
            df = t.df
            ws = wb.create_sheet(title=f"Table{i+1}")
            for r in dataframe_to_rows := __import__('openpyxl').utils.dataframe.dataframe_to_rows(df, index=False, header=True):
                ws.append(r)
        wb.save(out_xlsx)
        return FileResponse(out_xlsx, filename='tables.xlsx')
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Excel -> PDF using LibreOffice
@app.post('/convert/excel-to-pdf')
async def excel_to_pdf(file: UploadFile = File(...)):
    input_path = save_upload(file)
    outdir = os.path.dirname(input_path)
    cmd = f"soffice --headless --convert-to pdf --outdir {outdir} {input_path}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    pdfs = list(pathlib.Path(outdir).glob('*.pdf'))
    if pdfs:
        return FileResponse(str(sorted(pdfs, key=lambda p:p.stat().st_mtime)[-1]), filename=sorted(pdfs)[-1].name)
    raise HTTPException(status_code=500, detail='pdf not produced')

# PPT -> PDF using LibreOffice
@app.post('/convert/ppt-to-pdf')
async def ppt_to_pdf(file: UploadFile = File(...)):
    input_path = save_upload(file)
    outdir = os.path.dirname(input_path)
    cmd = f"soffice --headless --convert-to pdf --outdir {outdir} {input_path}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    pdfs = list(pathlib.Path(outdir).glob('*.pdf'))
    if pdfs:
        return FileResponse(str(sorted(pdfs, key=lambda p:p.stat().st_mtime)[-1]), filename=sorted(pdfs)[-1].name)
    raise HTTPException(status_code=500, detail='pdf not produced')

# PDF -> PPT (image based) using pdftoppm + python-pptx
@app.post('/convert/pdf-to-ppt')
async def pdf_to_ppt(file: UploadFile = File(...), dpi: int = Form(150)):
    input_path = save_upload(file)
    prefix = os.path.join(TMP, uuid.uuid4().hex + '_slide')
    cmd = f"pdftoppm -jpeg -r {dpi} {input_path} {prefix}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    imgs = sorted([str(p) for p in pathlib.Path(TMP).glob(os.path.basename(prefix)+'*')])
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        for img in imgs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
        out_ppt = make_out_path('.pptx')
        prs.save(out_ppt)
        return FileResponse(out_ppt, filename='presentation.pptx')
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Repair PDF using Ghostscript
@app.post('/tools/repair')
async def repair_pdf(file: UploadFile = File(...)):
    input_path = save_upload(file)
    out = make_out_path('_repaired.pdf')
    cmd = f"gs -o {out} -sDEVICE=pdfwrite -dPDFSETTINGS=/prepress {input_path}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out, filename='repaired.pdf')

# Convert to PDF/A (basic)
@app.post('/tools/pdfa')
async def convert_pdfa(file: UploadFile = File(...)):
    input_path = save_upload(file)
    out = make_out_path('_pdfa.pdf')
    # This is a simplified PDF/A conversion command (may require ICC profiles in production)
    cmd = f"gs -dPDFA -dBATCH -dNOPAUSE -sProcessColorModel=DeviceRGB -sDEVICE=pdfwrite -sOutputFile={out} {input_path}"
    rc,o,e = run_cmd(cmd)
    if rc!=0:
        raise HTTPException(status_code=500, detail={'cmd':cmd,'out':o,'err':e})
    return FileResponse(out, filename='pdfa.pdf')

# Edit PDF - add text overlay at coordinates using PyMuPDF (fitz)
@app.post('/tools/edit/add-text')
async def edit_add_text(file: UploadFile = File(...), page: int = Form(1), x: float = Form(50.0), y: float = Form(50.0), text: str = Form(...), fontsize: int = Form(12)):
    try:
        import fitz  # PyMuPDF
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing PyMuPDF: {e}")
    input_path = save_upload(file)
    doc = fitz.open(input_path)
    p = doc[page-1]
    p.insert_text((x, y), text, fontsize=fontsize)
    out = make_out_path('_edited.pdf')
    doc.save(out)
    doc.close()
    return FileResponse(out, filename='edited.pdf')

# HTML to PDF using WeasyPrint (if available)
@app.post('/convert/html-to-pdf')
async def html_to_pdf(html: str = Form(...)):
    try:
        from weasyprint import HTML
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Missing WeasyPrint: {e}")
    out = make_out_path('.pdf')
    HTML(string=html).write_pdf(out)
    return FileResponse(out, filename='page.pdf')
